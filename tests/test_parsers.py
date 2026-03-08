import io
import pytest


def make_pdf_bytes(page_texts: list[str]) -> bytes:
    from pypdf import PdfWriter
    writer = PdfWriter()
    for text in page_texts:
        page = writer.add_blank_page(width=612, height=792)
        # pypdf blank pages have no text; we test the parser logic via a workaround
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def test_pptx_parser_slide_numbers():
    from pptx import Presentation
    from pptx.util import Inches
    from rag.ingestion.pptx_parser import parse_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank
    for i in range(3):
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = f"Content of slide {i + 1}"

    buf = io.BytesIO()
    prs.save(buf)
    chunks = parse_pptx(buf.getvalue(), "test.pptx")

    assert len(chunks) == 3
    assert chunks[0]["slide"] == 1
    assert chunks[1]["slide"] == 2
    assert chunks[2]["slide"] == 3
    assert "Content of slide 1" in chunks[0]["text"]
    assert all(c["source"] == "test.pptx" for c in chunks)


def test_docx_parser_basic():
    from docx import Document
    from rag.ingestion.docx_parser import parse_docx

    doc = Document()
    for i in range(5):
        doc.add_paragraph(f"Paragraph {i + 1}: " + "word " * 50)

    buf = io.BytesIO()
    doc.save(buf)
    chunks = parse_docx(buf.getvalue(), "test.docx")

    assert len(chunks) >= 1
    assert all(c["source"] == "test.docx" for c in chunks)
    assert all("page" in c for c in chunks)


def test_loader_dispatch_unknown_raises():
    from rag.ingestion.loader import load_and_chunk
    with pytest.raises(ValueError, match="Unsupported"):
        load_and_chunk(b"data", "file.csv")
