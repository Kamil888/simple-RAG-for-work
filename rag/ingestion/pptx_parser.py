import io
from pptx import Presentation

CHUNK_SIZE = 2000
CHUNK_OVERLAP = 200


def _split_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return [c for c in chunks if c.strip()]


def parse_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []
    chunk_index = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)

        slide_text = "\n".join(texts).strip()
        if not slide_text:
            continue

        for fragment in _split_text(slide_text):
            chunks.append({
                "text": fragment,
                "source": filename,
                "slide": slide_num,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

    return chunks
